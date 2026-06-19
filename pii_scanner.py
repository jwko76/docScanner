#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PiiScanner.py  ─  개인정보/민감정보 파일 스캐너 (Python 포터블 버전)
====================================================================
사용법:
  python pii_scanner.py [옵션]
  python pii_scanner.py --path "C:\\Users\\홍길동\\Documents" --output "C:\\Reports"
  python pii_scanner.py --skip-images --threads 8

기본 동작 (보안 설계):
  * 기본 스캔 경로: 현재 사용자 폴더(%USERPROFILE%)만 대상
  * Windows / Program Files / ProgramData 등 시스템 폴더 항상 제외
  * 레지스트리 미접근 (winreg 사용 없음)
  * 스캔 중 네트워크 미접근 (파일 로컬 읽기만 수행)
  * 결과 파일은 --output 지정 폴더에만 저장
  * 전체 드라이브 스캔이 필요하면 --all-drives 명시

의존성 설치:
  pip install -r requirements.txt

OCR 사용 시:
  Tesseract 설치 필요 → https://github.com/UB-Mannheim/tesseract/wiki
  설치 시 Korean 언어팩 선택
"""

import argparse
import ctypes
import ctypes.wintypes
import datetime
import json
import os
import re
import sys
import threading
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Dict, Tuple

# ──────────────────────────────────────────────
# 의존성 자동 설치
# ──────────────────────────────────────────────
def _auto_install(pkg: str):
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", pkg, "--quiet"],
                          stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

REQUIRED = {
    "docx":       "python-docx",
    "openpyxl":   "openpyxl",
    "pptx":       "python-pptx",
    "pdfplumber": "pdfplumber",
    "xlsxwriter": "xlsxwriter",
    "tqdm":       "tqdm",
    "chardet":    "chardet",
    "colorama":   "colorama",
    "PIL":        "Pillow",
}

print("의존성 확인 중...", end="", flush=True)
for mod, pkg in REQUIRED.items():
    try:
        __import__(mod)
    except ImportError:
        print(f"\n  {pkg} 설치 중...", end="", flush=True)
        try:
            _auto_install(pkg)
        except Exception:
            print(f" [실패 - 수동 설치 필요: pip install {pkg}]")
print(" OK")

# ──────────────────────────────────────────────
# 이제 import
# ──────────────────────────────────────────────
import chardet
import xlsxwriter
from colorama import Fore, Style, init as colorama_init
from tqdm import tqdm

colorama_init(autoreset=True)

# 선택적 import
try:
    import docx as python_docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    import pptx as python_pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import pdfplumber
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import pytesseract
    from PIL import Image
    HAS_OCR = True
    # Tesseract 실행 파일 확인
    try:
        pytesseract.get_tesseract_version()
    except Exception:
        HAS_OCR = False
except ImportError:
    HAS_OCR = False

try:
    import olefile
    HAS_OLE = True
except ImportError:
    HAS_OLE = False

# ════════════════════════════════════════════════════════════════════
# 1. Everything SDK 스캐너
# ════════════════════════════════════════════════════════════════════

DOCUMENT_EXTENSIONS = {
    ".hwp", ".hwpx",
    ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".pdf", ".txt", ".log", ".csv", ".tsv", ".rtf",
    ".xml", ".json", ".html", ".htm",
    ".ini", ".cfg", ".conf",
    ".odt", ".ods", ".odp",
}
IMAGE_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".tif", ".tiff",
    ".bmp", ".gif", ".webp",
}

# ──────────────────────────────────────────────────────────────────────
# 시스템 폴더 제외 목록
#   - Windows 핵심 디렉터리 / 프로그램 설치 폴더 / 가상드라이브 등
#   - 개인정보가 존재할 이유가 없고, 스캔해도 의미 없는 경로
# ──────────────────────────────────────────────────────────────────────
_SYS_DRIVE = os.environ.get("SystemDrive", "C:").rstrip("\\")

SYSTEM_EXCLUDED_PATHS: List[str] = [
    # Windows 핵심
    _SYS_DRIVE + "\\Windows",
    _SYS_DRIVE + "\\Windows.old",
    # 프로그램 설치
    _SYS_DRIVE + "\\Program Files",
    _SYS_DRIVE + "\\Program Files (x86)",
    _SYS_DRIVE + "\\ProgramData",
    # 복구·시스템 영역
    _SYS_DRIVE + "\\Recovery",
    _SYS_DRIVE + "\\$Recycle.Bin",
    _SYS_DRIVE + "\\$WINDOWS.~BT",
    _SYS_DRIVE + "\\System Volume Information",
    _SYS_DRIVE + "\\Boot",
    _SYS_DRIVE + "\\bootmgr",
    # WinSxS / 패키지 캐시
    _SYS_DRIVE + "\\Windows\\SysWOW64",
    _SYS_DRIVE + "\\Windows\\System32",
    # 가상화
    _SYS_DRIVE + "\\ProgramData\\Microsoft\\Windows\\Hyper-V",
]

def _is_excluded(path: str) -> bool:
    """경로가 시스템 제외 목록에 해당하면 True"""
    path_lower = path.lower()
    for ex in SYSTEM_EXCLUDED_PATHS:
        if path_lower.startswith(ex.lower()):
            return True
    return False


def get_user_default_paths() -> List[str]:
    """
    현재 로그인 사용자의 기본 데이터 폴더 목록을 반환.
    개인정보가 저장될 가능성이 높은 경로만 포함.
    레지스트리 미사용 — 환경변수로만 경로 조회.
    """
    profile = os.environ.get("USERPROFILE", "")
    appdata  = os.environ.get("APPDATA", "")        # %APPDATA% = Roaming
    localapp = os.environ.get("LOCALAPPDATA", "")   # LocalAppData

    candidates = [
        profile,                                    # 전체 홈 (Documents/Desktop/Downloads 포함)
        appdata,                                    # 앱 사용자 설정 (Roaming)
        localapp,                                   # 앱 사용자 설정 (Local)
    ]
    # OneDrive 경로 (설치된 경우)
    for env in ("OneDrive", "OneDriveConsumer", "OneDriveCommercial"):
        od = os.environ.get(env, "")
        if od and os.path.isdir(od):
            candidates.append(od)

    # 실제로 존재하는 폴더만 반환
    return [p for p in candidates if p and os.path.isdir(p)]

@dataclass
class FileEntry:
    path: str
    extension: str
    size: int
    is_document: bool
    is_image: bool


class EverythingScanner:
    """Everything SDK DLL을 ctypes로 호출하여 파일 목록 조회"""

    # Everything 상수
    SORT_PATH_ASCENDING = 3
    REQUEST_FULL_PATH   = 0x00000004
    REQUEST_SIZE        = 0x00000010
    REQUEST_DATE_MOD    = 0x00000040

    def __init__(self):
        self._dll = None
        self._lock = threading.Lock()

    def _find_dll(self) -> Optional[str]:
        candidates = [
            Path(__file__).parent / "Everything64.dll",
            Path(__file__).parent / "sdk" / "Everything64.dll",
            Path(__file__).parent / "Everything32.dll",
        ]
        for c in candidates:
            if c.exists():
                return str(c)
        # PATH에서 탐색
        for d in os.environ.get("PATH", "").split(";"):
            for name in ("Everything64.dll", "Everything32.dll"):
                p = Path(d) / name
                if p.exists():
                    return str(p)
        return None

    def initialize(self, dll_path: str = "") -> bool:
        path = dll_path or self._find_dll()
        if not path:
            print(f"{Fore.RED}[오류] Everything64.dll을 찾을 수 없습니다.")
            print("  → exe와 같은 폴더에 복사하거나 --dll 옵션으로 지정하세요.")
            return False
        try:
            self._dll = ctypes.WinDLL(path)
            self._setup_functions()
            if not self._dll.Everything_IsDBLoaded():
                print(f"{Fore.RED}[오류] Everything DB가 로드되지 않았습니다.")
                print("  → Everything을 실행한 후 재시도하세요.")
                return False
            return True
        except Exception as e:
            print(f"{Fore.RED}[오류] DLL 로드 실패: {e}")
            return False

    def _setup_functions(self):
        d = self._dll
        d.Everything_SetSearchW.argtypes       = [ctypes.c_wchar_p]
        d.Everything_SetSearchW.restype        = None
        d.Everything_SetMax.argtypes           = [ctypes.wintypes.DWORD]
        d.Everything_SetMax.restype            = None
        d.Everything_SetOffset.argtypes        = [ctypes.wintypes.DWORD]
        d.Everything_SetOffset.restype         = None
        d.Everything_SetSort.argtypes          = [ctypes.wintypes.DWORD]
        d.Everything_SetSort.restype           = None
        d.Everything_SetRequestFlags.argtypes  = [ctypes.wintypes.DWORD]
        d.Everything_SetRequestFlags.restype   = None
        d.Everything_SetMatchCase.argtypes     = [ctypes.wintypes.BOOL]
        d.Everything_SetMatchCase.restype      = None
        d.Everything_QueryW.argtypes           = [ctypes.wintypes.BOOL]
        d.Everything_QueryW.restype            = ctypes.wintypes.BOOL
        d.Everything_GetNumResults.argtypes    = []
        d.Everything_GetNumResults.restype     = ctypes.wintypes.DWORD
        d.Everything_IsFileResult.argtypes     = [ctypes.wintypes.DWORD]
        d.Everything_IsFileResult.restype      = ctypes.wintypes.BOOL
        d.Everything_GetResultFullPathNameW.argtypes = [
            ctypes.wintypes.DWORD, ctypes.c_wchar_p, ctypes.wintypes.DWORD]
        d.Everything_GetResultFullPathNameW.restype = None
        d.Everything_GetResultSize.argtypes    = [
            ctypes.wintypes.DWORD, ctypes.POINTER(ctypes.c_int64)]
        d.Everything_GetResultSize.restype     = ctypes.wintypes.BOOL
        d.Everything_IsDBLoaded.argtypes       = []
        d.Everything_IsDBLoaded.restype        = ctypes.wintypes.BOOL
        d.Everything_Reset.argtypes            = []
        d.Everything_Reset.restype             = None
        d.Everything_CleanUp.argtypes          = []
        d.Everything_CleanUp.restype           = None

    def _build_query(self, root_paths: List[str]) -> str:
        """
        root_paths: 스캔할 경로 목록.
        비어있으면 전체 드라이브 (--all-drives 전용).
        여러 경로면 OR로 묶어 쿼리.
        """
        all_exts = DOCUMENT_EXTENSIONS | IMAGE_EXTENSIONS
        ext_parts = [f"ext:{e.lstrip('.')}" for e in all_exts]
        ext_query = " | ".join(ext_parts)
        if not root_paths:
            return f"({ext_query})"
        # 각 경로를 "path" 형태로 OR 결합
        path_parts = " | ".join(f'"{p}"' for p in root_paths)
        return f"({path_parts}) ({ext_query})"

    def scan(self, root_paths: List[str], exclude_system: bool = True,
             extra_excludes: List[str] = None, progress_cb=None) -> List[FileEntry]:
        """
        root_paths: 스캔 경로 목록 (빈 리스트 = 전체 드라이브)
        exclude_system: True이면 시스템 폴더를 결과에서 항상 제외
        extra_excludes: 추가로 제외할 경로 목록
        """
        if not self._dll:
            return []

        query = self._build_query(root_paths)
        d = self._dll

        d.Everything_Reset()
        d.Everything_SetSearchW(query)
        d.Everything_SetMatchCase(False)
        d.Everything_SetSort(self.SORT_PATH_ASCENDING)
        d.Everything_SetRequestFlags(
            self.REQUEST_FULL_PATH | self.REQUEST_SIZE | self.REQUEST_DATE_MOD)
        d.Everything_SetMax(1_000_000)
        d.Everything_SetOffset(0)

        print(f"  쿼리 실행 중...", end="", flush=True)
        if not d.Everything_QueryW(True):
            print(f" {Fore.RED}실패")
            return []

        total = d.Everything_GetNumResults()
        print(f" {total:,}개 파일 발견")

        # 추가 제외 경로 정규화
        extra_lower = [e.lower() for e in (extra_excludes or [])]

        buf = ctypes.create_unicode_buffer(32768)
        results = []
        skipped_sys = 0

        for i in range(total):
            if not d.Everything_IsFileResult(i):
                continue
            d.Everything_GetResultFullPathNameW(i, buf, 32768)
            path = buf.value
            path_lower = path.lower()

            # ── 시스템 경로 필터링 ──────────────────────────────────
            if exclude_system and _is_excluded(path):
                skipped_sys += 1
                continue
            if any(path_lower.startswith(ex) for ex in extra_lower):
                skipped_sys += 1
                continue

            ext = Path(path).suffix.lower()
            size_val = ctypes.c_int64(0)
            d.Everything_GetResultSize(i, ctypes.byref(size_val))

            entry = FileEntry(
                path=path,
                extension=ext,
                size=size_val.value,
                is_document=(ext in DOCUMENT_EXTENSIONS),
                is_image=(ext in IMAGE_EXTENSIONS),
            )
            results.append(entry)

            if progress_cb and i % 5000 == 0:
                progress_cb(i, total)

        d.Everything_CleanUp()
        if skipped_sys:
            print(f"  (시스템 폴더 제외: {skipped_sys:,}개 건너뜀)")
        return results


# ════════════════════════════════════════════════════════════════════
# 2. 텍스트 추출기
# ════════════════════════════════════════════════════════════════════

MAX_TEXT_LEN = 2_000_000   # 최대 2MB 텍스트


def _safe_str(s) -> str:
    if s is None:
        return ""
    if isinstance(s, bytes):
        return s.decode("utf-8", errors="replace")
    return str(s)


def extract_text(file_path: str, extension: str, use_ocr: bool = True) -> Tuple[str, str, str]:
    """
    Returns: (text, method, error)
    """
    ext = extension.lower()

    try:
        # ── 텍스트 계열 ──────────────────────────────────
        if ext in {".txt", ".log", ".csv", ".tsv", ".ini", ".cfg", ".conf",
                   ".xml", ".json", ".html", ".htm"}:
            return _extract_text_file(file_path)

        # ── DOCX ──────────────────────────────────────────
        if ext == ".docx":
            return _extract_docx(file_path)

        # ── XLSX ──────────────────────────────────────────
        if ext == ".xlsx":
            return _extract_xlsx(file_path)

        # ── PPTX ──────────────────────────────────────────
        if ext == ".pptx":
            return _extract_pptx(file_path)

        # ── PDF ───────────────────────────────────────────
        if ext == ".pdf":
            return _extract_pdf(file_path)

        # ── HWP ───────────────────────────────────────────
        if ext in {".hwp", ".hwpx"}:
            return _extract_hwp(file_path, ext)

        # ── 구버전 OLE (doc / xls / ppt) ─────────────────
        if ext in {".doc", ".xls", ".ppt"}:
            return _extract_ole(file_path, ext)

        # ── RTF ───────────────────────────────────────────
        if ext == ".rtf":
            return _extract_rtf(file_path)

        # ── 이미지 OCR ────────────────────────────────────
        if ext in IMAGE_EXTENSIONS:
            if not use_ocr:
                return "", "skipped", "OCR 건너뜀"
            return _extract_image_ocr(file_path)

        return "", "unknown", f"지원하지 않는 확장자: {ext}"

    except PermissionError:
        return "", "error", "접근 권한 없음"
    except FileNotFoundError:
        return "", "error", "파일 없음"
    except Exception as e:
        return "", "error", f"추출 오류: {type(e).__name__}: {str(e)[:100]}"


def _detect_encoding(data: bytes) -> str:
    if data[:3] == b"\xef\xbb\xbf":
        return "utf-8-sig"
    if data[:2] == b"\xff\xfe":
        return "utf-16-le"
    if data[:2] == b"\xfe\xff":
        return "utf-16-be"
    detected = chardet.detect(data[:8192])
    enc = detected.get("encoding") or "utf-8"
    # EUC-KR 별칭 정규화
    if enc.lower() in ("euc-kr", "euc_kr", "cp949", "ms949"):
        return "cp949"
    return enc


def _extract_text_file(path: str) -> Tuple[str, str, str]:
    with open(path, "rb") as f:
        raw = f.read(MAX_TEXT_LEN * 3)
    enc = _detect_encoding(raw)
    try:
        text = raw.decode(enc, errors="replace")
    except Exception:
        text = raw.decode("utf-8", errors="replace")
    return text[:MAX_TEXT_LEN], "PlainText", ""


def _extract_docx(path: str) -> Tuple[str, str, str]:
    if not HAS_DOCX:
        return "", "error", "python-docx 미설치"
    doc = python_docx.Document(path)
    parts = []
    for para in doc.paragraphs:
        parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    text = "\n".join(parts)
    return text[:MAX_TEXT_LEN], "python-docx", ""


def _extract_xlsx(path: str) -> Tuple[str, str, str]:
    if not HAS_OPENPYXL:
        return "", "error", "openpyxl 미설치"
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    parts.append(str(cell))
    wb.close()
    text = " ".join(parts)
    return text[:MAX_TEXT_LEN], "openpyxl", ""


def _extract_pptx(path: str) -> Tuple[str, str, str]:
    if not HAS_PPTX:
        return "", "error", "python-pptx 미설치"
    prs = python_pptx.Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    text = "\n".join(parts)
    return text[:MAX_TEXT_LEN], "python-pptx", ""


def _extract_pdf(path: str) -> Tuple[str, str, str]:
    if not HAS_PDF:
        return "", "error", "pdfplumber 미설치"
    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                parts.append(t)
            if sum(len(p) for p in parts) > MAX_TEXT_LEN:
                break
    text = "\n".join(parts)
    return text[:MAX_TEXT_LEN], "pdfplumber", ""


def _extract_hwp(path: str, ext: str) -> Tuple[str, str, str]:
    # hwp5 라이브러리 시도
    try:
        import hwp5.binmodel  # type: ignore
        from hwp5.hwp5txt import HwpDocument  # type: ignore
        doc = HwpDocument(path)
        text = doc.get_text()
        return text[:MAX_TEXT_LEN], "hwp5", ""
    except ImportError:
        pass
    except Exception as e:
        pass

    # hwpx는 OOXML 유사 ZIP 구조
    if ext == ".hwpx":
        return _extract_hwpx_zip(path)

    # fallback: Windows COM (HWP 설치 필요)
    return _extract_via_com_hwp(path)


def _extract_hwpx_zip(path: str) -> Tuple[str, str, str]:
    """hwpx는 ZIP+XML 구조"""
    import zipfile
    parts = []
    try:
        with zipfile.ZipFile(path, "r") as zf:
            names = zf.namelist()
            xml_files = [n for n in names if n.endswith(".xml")]
            for xf in xml_files[:20]:
                try:
                    data = zf.read(xf).decode("utf-8", errors="replace")
                    # XML 태그 제거
                    clean = re.sub(r"<[^>]+>", " ", data)
                    parts.append(clean)
                except Exception:
                    pass
        text = " ".join(parts)
        return text[:MAX_TEXT_LEN], "hwpx-zip", ""
    except Exception as e:
        return "", "error", f"hwpx 파싱 실패: {e}"


def _extract_via_com_hwp(path: str) -> Tuple[str, str, str]:
    """Windows COM으로 HWP 추출 (한컴 오피스 설치 필요)"""
    try:
        import win32com.client  # type: ignore
        hwp = win32com.client.Dispatch("HWPFrame.HwpObject")
        hwp.RegisterModule("FilePathCheckDLL", "FilePathCheckerModuleExample")
        hwp.Open(str(Path(path).resolve()))
        text = hwp.GetTextFile("TEXT", "")
        hwp.Quit()
        return text[:MAX_TEXT_LEN], "COM-HWP", ""
    except Exception as e:
        return "", "error", f"HWP COM 실패 (한컴오피스 미설치?): {str(e)[:80]}"


def _extract_ole(path: str, ext: str) -> Tuple[str, str, str]:
    """구버전 Office 파일 (.doc/.xls/.ppt) - olefile로 텍스트 추출"""
    if not HAS_OLE:
        return "", "error", "olefile 미설치"
    try:
        # WordDocument 스트림에서 텍스트 추출 (간단 방식)
        ole = olefile.OleFileIO(path)
        parts = []
        for entry in ole.listdir():
            name = "/".join(entry)
            if any(k in name for k in ("WordDocument", "Workbook", "PowerPoint")):
                try:
                    data = ole.openstream(entry).read()
                    # 단순 UTF-16-LE 텍스트 추출
                    decoded = data.decode("utf-16-le", errors="ignore")
                    printable = "".join(c for c in decoded if c.isprintable() or c in "\n\t ")
                    parts.append(printable)
                except Exception:
                    pass
        ole.close()
        text = " ".join(parts)
        if len(text.strip()) < 50:
            # fallback: antiword / COM
            return _extract_ole_com(path, ext)
        return text[:MAX_TEXT_LEN], "olefile", ""
    except Exception as e:
        return _extract_ole_com(path, ext)


def _extract_ole_com(path: str, ext: str) -> Tuple[str, str, str]:
    """Windows COM으로 구버전 Office 파일 추출"""
    try:
        import win32com.client  # type: ignore
        abs_path = str(Path(path).resolve())
        if ext == ".doc":
            app = win32com.client.Dispatch("Word.Application")
            app.Visible = False
            doc = app.Documents.Open(abs_path, ReadOnly=True)
            text = doc.Content.Text
            doc.Close(False)
            app.Quit()
            return text[:MAX_TEXT_LEN], "COM-Word", ""
        elif ext == ".xls":
            app = win32com.client.Dispatch("Excel.Application")
            app.Visible = False
            wb = app.Workbooks.Open(abs_path)
            parts = []
            for ws in wb.Worksheets:
                used = ws.UsedRange
                parts.append(str(used.Value))
            wb.Close(False)
            app.Quit()
            return " ".join(parts)[:MAX_TEXT_LEN], "COM-Excel", ""
    except Exception as e:
        return "", "error", f"COM Office 실패: {str(e)[:80]}"
    return "", "error", "구버전 Office 추출 실패"


def _extract_rtf(path: str) -> Tuple[str, str, str]:
    """RTF에서 텍스트 추출 (정규식으로 태그 제거)"""
    with open(path, "rb") as f:
        raw = f.read(MAX_TEXT_LEN * 2)
    enc = _detect_encoding(raw)
    text = raw.decode(enc, errors="replace")
    # RTF 태그 제거
    text = re.sub(r"\\[a-z]+\d*\s?", " ", text)
    text = re.sub(r"[{}\\]", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text[:MAX_TEXT_LEN], "RTF-regex", ""


def _extract_image_ocr(path: str) -> Tuple[str, str, str]:
    """Tesseract OCR로 이미지 텍스트 추출"""
    if not HAS_OCR:
        return "", "error", "Tesseract 미설치 (https://github.com/UB-Mannheim/tesseract/wiki)"
    try:
        img = Image.open(path)
        # 해상도 향상
        if max(img.size) < 1000:
            scale = 1000 / max(img.size)
            new_size = (int(img.width * scale), int(img.height * scale))
            img = img.resize(new_size, Image.LANCZOS)

        # 한국어 + 영어 OCR
        lang = "kor+eng" if "kor" in pytesseract.get_languages() else "eng"
        text = pytesseract.image_to_string(img, lang=lang)
        return text[:MAX_TEXT_LEN], f"Tesseract({lang})", ""
    except Exception as e:
        return "", "error", f"OCR 실패: {str(e)[:100]}"


# ════════════════════════════════════════════════════════════════════
# 3. 개인정보 탐지기
# ════════════════════════════════════════════════════════════════════

@dataclass
class PiiMatch:
    pii_type: str
    matched: str
    masked: str
    line_no: int
    context: str
    confidence: float


@dataclass
class FileScanResult:
    path: str
    extension: str
    success: bool
    error: str
    method: str
    text_len: int
    scan_sec: float
    matches: List[PiiMatch] = field(default_factory=list)

    @property
    def match_count(self):
        return len(self.matches)


class PiiDetector:
    """정규식 + 검증 로직으로 개인정보 탐지"""

    # ── 패턴 정의 ──────────────────────────────────────────────────
    PATTERNS = [
        # 주민등록번호: YYMMDD-NNNNNNN (체크섬 검증)
        ("주민등록번호", re.compile(
            r"\b(\d{6})([-\s])([1-489]\d{6})\b", re.UNICODE
        ), True),

        # 전화번호 (국내)
        ("전화번호", re.compile(
            r"\b0(?:1[016-9]|2|3[1-3]|4[1-4]|5[1-5]|6[1-4]|7[1-4])"
            r"[\s\-\.]?\d{3,4}[\s\-\.]?\d{4}\b"
            r"|\b1[5-8]\d{2}[\s\-\.]?\d{4}\b",
            re.UNICODE
        ), False),

        # 이메일
        ("이메일", re.compile(
            r"\b[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}\b"
        ), False),

        # IPv4
        ("IP주소", re.compile(
            r"\b(?:(?:25[0-5]|2[0-4]\d|[01]?\d\d?)\.){3}"
            r"(?:25[0-5]|2[0-4]\d|[01]?\d\d?)\b"
        ), True),

        # MAC 주소
        ("MAC주소", re.compile(
            r"\b(?:[0-9A-Fa-f]{2}[:\-]){5}[0-9A-Fa-f]{2}\b"
        ), False),

        # 신용카드 (Luhn 검증)
        ("신용카드", re.compile(
            r"\b(?:4[0-9]{3}|5[1-5][0-9]{2}|3[47][0-9]{2}|6(?:011|5[0-9]{2}))"
            r"[\s\-]?\d{4}[\s\-]?\d{4}[\s\-]?\d{4}\b"
            r"|\b\d{4}[\s\-]\d{4}[\s\-]\d{4}[\s\-]\d{4}\b",
            re.UNICODE
        ), True),

        # 계좌번호 (국내 은행 패턴)
        ("계좌번호", re.compile(
            r"\b\d{3,4}[-\s]\d{4,6}[-\s]\d{2,7}\b"
        ), False),

        # 여권번호: 한국 여권 M+8자리 또는 영문 2자리+7자리
        ("여권번호", re.compile(r"\b[A-Z]{1,2}\d{7,8}\b"
        ), False),

        # 운전면허번호
        ("운전면허", re.compile(
            r"\b(?:0[1-9]|[1-2]\d|3[0-6])-\d{2}-\d{6}-\d{2}\b"
        ), False),
    ]

    # 한국 주소 패턴 (광역시/도 키워드 기반)
    ADDR_PATTERN = re.compile(
        r"(?:서울|부산|대구|인천|광주|대전|울산|세종|경기|강원|충북|충남|전북|전남|경북|경남|제주)"
        r".{0,30}?"
        r"(?:시|도|구|군|동|읍|면|로|길)\s*\d+",
        re.UNICODE | re.DOTALL
    )

    CONTEXT_LEN = 80

    def detect(self, text: str) -> List[PiiMatch]:
        if not text:
            return []
        results = []
        lines = text.split("\n")
        # 줄 번호 빠른 조회를 위한 오프셋 테이블
        line_starts = []
        pos = 0
        for ln in lines:
            line_starts.append(pos)
            pos += len(ln) + 1

        def get_line(offset: int) -> int:
            lo, hi = 0, len(line_starts) - 1
            while lo < hi:
                mid = (lo + hi + 1) // 2
                if line_starts[mid] <= offset:
                    lo = mid
                else:
                    hi = mid - 1
            return lo + 1  # 1-based

        def get_context(offset: int, length: int) -> str:
            half = self.CONTEXT_LEN // 2
            start = max(0, offset - half)
            end = min(len(text), offset + length + half)
            ctx = text[start:end].replace("\n", " ").replace("\r", "")
            return ctx

        def mask(pii_type: str, matched: str) -> str:
            if pii_type == "주민등록번호":
                parts = re.split(r"[-\s]", matched)
                if len(parts) == 2:
                    return parts[0] + "-" + "*" * 7
                return matched[:6] + "-*******"
            if pii_type == "신용카드":
                digits = re.sub(r"\D", "", matched)
                return "****-****-****-" + digits[-4:] if len(digits) >= 4 else "****"
            if pii_type == "이메일":
                at = matched.find("@")
                if at > 1:
                    return matched[0] + "*" * (at - 1) + matched[at:]
                return "***@" + matched.split("@")[-1]
            if pii_type == "전화번호":
                return re.sub(r"(\d{3,4})([-\s])(\d{3,4})([-\s])(\d{4})",
                              lambda m: m.group(1)+m.group(2)+"****"+m.group(4)+m.group(5),
                              matched) if matched != matched else matched[:4] + "****" + matched[-4:]
            # 기본: 절반 마스킹
            half = len(matched) // 2
            return matched[:half] + "*" * (len(matched) - half)

        # 정규식 패턴 탐지
        for type_name, pattern, do_validate in self.PATTERNS:
            for m in pattern.finditer(text):
                matched = m.group()
                offset  = m.start()
                # 검증
                confidence = 0.8
                if do_validate:
                    if type_name == "주민등록번호":
                        if not self._validate_rrn(matched):
                            continue
                        confidence = 0.99
                    elif type_name == "신용카드":
                        if not self._luhn(matched):
                            continue
                        confidence = 0.95
                    elif type_name == "IP주소":
                        if not self._validate_ip(matched):
                            continue
                        confidence = 0.85

                results.append(PiiMatch(
                    pii_type=type_name,
                    matched=matched,
                    masked=mask(type_name, matched),
                    line_no=get_line(offset),
                    context=get_context(offset, len(matched)),
                    confidence=confidence,
                ))
                if len(results) > 10_000:
                    return results

        # 주소 탐지
        for m in self.ADDR_PATTERN.finditer(text):
            matched = m.group().strip()
            if len(matched) < 8:
                continue
            offset = m.start()
            results.append(PiiMatch(
                pii_type="주소",
                matched=matched[:80],
                masked=matched[:len(matched)//2] + "*" * (len(matched) - len(matched)//2),
                line_no=get_line(offset),
                context=get_context(offset, len(matched)),
                confidence=0.7,
            ))

        return results

    # ── 검증 함수 ──────────────────────────────────────────────────

    @staticmethod
    def _validate_rrn(text: str) -> bool:
        digits = re.sub(r"\D", "", text)
        if len(digits) != 13:
            return False
        # 날짜 유효성
        month = int(digits[2:4])
        day   = int(digits[4:6])
        if not (1 <= month <= 12 and 1 <= day <= 31):
            return False
        # 체크섬
        weights = [2, 3, 4, 5, 6, 7, 8, 9, 2, 3, 4, 5]
        total = sum(int(digits[i]) * weights[i] for i in range(12))
        check = (11 - total % 11) % 10
        return check == int(digits[12])

    @staticmethod
    def _luhn(text: str) -> bool:
        digits = re.sub(r"\D", "", text)
        if len(digits) < 13:
            return False
        total = 0
        for i, d in enumerate(reversed(digits)):
            n = int(d)
            if i % 2 == 1:
                n *= 2
                if n > 9:
                    n -= 9
            total += n
        return total % 10 == 0

    @staticmethod
    def _validate_ip(text: str) -> bool:
        parts = text.split(".")
        if len(parts) != 4:
            return False
        vals = [int(p) for p in parts]
        # 0.0.0.0, 255.255.x.x 등 제외
        if vals[0] == 0 or all(v == 255 for v in vals[:2]):
            return False
        return True


# ════════════════════════════════════════════════════════════════════
# 4. 리포트 생성기
# ════════════════════════════════════════════════════════════════════

class Reporter:

    def save_all(self, results: List[FileScanResult], summary: dict,
                 output_dir: str, base_name: str):
        os.makedirs(output_dir, exist_ok=True)
        xlsx_path = os.path.join(output_dir, base_name + ".xlsx")
        html_path = os.path.join(output_dir, base_name + ".html")
        self._save_excel(results, summary, xlsx_path)
        self._save_html(results, summary, html_path)
        return xlsx_path, html_path

    # ── Excel ────────────────────────────────────────────────────

    def _save_excel(self, results: List[FileScanResult], summary: dict, path: str):
        wb = xlsxwriter.Workbook(path, {"strings_to_urls": False})

        # 서식
        hdr = wb.add_format({"bold": True, "bg_color": "#2F5496",
                              "font_color": "white", "border": 1, "align": "center"})
        cell = wb.add_format({"border": 1, "text_wrap": True})
        red  = wb.add_format({"border": 1, "font_color": "#C00000", "bold": True})
        num  = wb.add_format({"border": 1, "align": "right"})
        title= wb.add_format({"bold": True, "font_size": 14})

        # ── 요약 시트 ────────────────────────────────────────────
        ws = wb.add_worksheet("요약")
        ws.set_column(0, 0, 22)
        ws.set_column(1, 1, 35)
        ws.write(0, 0, "개인정보 스캔 결과 요약", title)
        rows = [
            ("스캔 경로",        summary["scan_path"]),
            ("스캔 일시",        summary["scan_time"]),
            ("스캔 파일 수",     summary["total_files"]),
            ("개인정보 파일 수", summary["files_with_pii"]),
            ("총 탐지 건수",     summary["total_pii"]),
            ("소요 시간",        f"{int(summary['elapsed'])//60}분 {int(summary['elapsed'])%60}초"),
        ]
        for i, (k, v) in enumerate(rows, 2):
            ws.write(i, 0, k, hdr)
            ws.write(i, 1, str(v), cell)

        ws.write(9, 0, "유형별 통계", title)
        ws.write(10, 0, "유형", hdr)
        ws.write(10, 1, "건수", hdr)
        for i, (t, c) in enumerate(summary["type_counts"].items(), 11):
            ws.write(i, 0, t, cell)
            ws.write(i, 1, c, num)

        # ── 파일 목록 시트 ───────────────────────────────────────
        wf = wb.add_worksheet("파일 목록")
        wf.set_column(0, 0, 55)
        wf.set_column(1, 5, 12)
        file_hdrs = ["파일 경로", "확장자", "탐지 건수", "추출 방법", "소요(초)", "오류"]
        for c, h in enumerate(file_hdrs):
            wf.write(0, c, h, hdr)
        for r, res in enumerate(results, 1):
            fmt = red if res.match_count > 0 else cell
            wf.write(r, 0, res.path,            cell)
            wf.write(r, 1, res.extension,        cell)
            wf.write(r, 2, res.match_count,      fmt)
            wf.write(r, 3, res.method,           cell)
            wf.write(r, 4, round(res.scan_sec,2), cell)
            wf.write(r, 5, res.error,            cell)
        wf.autofilter(0, 0, len(results), 5)

        # ── 상세 결과 시트 ───────────────────────────────────────
        wd = wb.add_worksheet("상세 결과")
        wd.set_column(0, 0, 50)
        wd.set_column(1, 1, 12)
        wd.set_column(2, 3, 22)
        wd.set_column(4, 4, 7)
        wd.set_column(5, 5, 60)
        det_hdrs = ["파일 경로", "탐지 유형", "탐지 값", "마스킹 값", "줄 번호", "맥락", "신뢰도"]
        for c, h in enumerate(det_hdrs):
            wd.write(0, c, h, hdr)
        dr = 1
        for res in results:
            for m in res.matches:
                wd.write(dr, 0, res.path,         cell)
                wd.write(dr, 1, m.pii_type,       red)
                wd.write(dr, 2, m.matched,        red)
                wd.write(dr, 3, m.masked,         cell)
                wd.write(dr, 4, m.line_no,        num)
                wd.write(dr, 5, m.context,        cell)
                wd.write(dr, 6, round(m.confidence, 2), num)
                dr += 1
        wd.autofilter(0, 0, dr, 6)

        wb.close()

    # ── HTML ─────────────────────────────────────────────────────

    def _save_html(self, results: List[FileScanResult], summary: dict, path: str):
        def he(s: str) -> str:
            return (str(s).replace("&","&amp;").replace("<","&lt;")
                    .replace(">","&gt;").replace('"',"&quot;"))

        type_counts = summary["type_counts"]
        max_count = max(type_counts.values(), default=1)

        badge_cls = {
            "주민등록번호": "ssn", "전화번호": "phone", "이메일": "email",
            "IP주소": "ip", "MAC주소": "mac", "신용카드": "card",
            "주소": "addr", "계좌번호": "bank",
            "여권번호": "other", "운전면허": "other",
        }

        html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>개인정보 스캔 결과</title>
<style>
body{{font-family:'맑은 고딕',sans-serif;margin:0;background:#f4f6f9;color:#333}}
.header{{background:linear-gradient(135deg,#1e3a5f,#2f5496);color:#fff;padding:24px 32px}}
.header h1{{margin:0;font-size:22px}} .header p{{margin:4px 0 0;opacity:.8;font-size:12px}}
.container{{max-width:1400px;margin:20px auto;padding:0 16px}}
.cards{{display:grid;grid-template-columns:repeat(auto-fit,minmax(160px,1fr));gap:14px;margin-bottom:20px}}
.card{{background:#fff;border-radius:8px;padding:14px 18px;box-shadow:0 2px 8px rgba(0,0,0,.08)}}
.card .val{{font-size:26px;font-weight:700;color:#2f5496}} .card .lbl{{font-size:11px;color:#888;margin-top:3px}}
.card.red .val{{color:#c00}}
.section{{background:#fff;border-radius:8px;box-shadow:0 2px 8px rgba(0,0,0,.08);margin-bottom:20px;overflow:hidden}}
.sh{{background:#2f5496;color:#fff;padding:10px 18px;font-weight:600;font-size:14px}}
table{{width:100%;border-collapse:collapse;font-size:12px}}
th{{background:#dce6f1;padding:7px 10px;text-align:left;font-weight:600;position:sticky;top:0}}
td{{padding:6px 10px;border-bottom:1px solid #eee;vertical-align:top;word-break:break-all}}
tr:hover td{{background:#f0f4fa}}
.badge{{display:inline-block;padding:1px 7px;border-radius:10px;font-size:11px;font-weight:600}}
.b-ssn{{background:#fde;color:#c00}} .b-phone{{background:#e9f0fd;color:#1a5276}}
.b-email{{background:#e9fde9;color:#196f3d}} .b-ip{{background:#fdf5e9;color:#784212}}
.b-mac{{background:#f5e9fd;color:#512e5f}} .b-card{{background:#fef9e7;color:#7d6608}}
.b-addr{{background:#e9fdfd;color:#0e6655}} .b-bank{{background:#fdecea;color:#922b21}}
.b-other{{background:#eee;color:#555}}
.red{{color:#c00;font-weight:600}} .gray{{color:#bbb}}
.bar-row{{display:flex;align-items:center;gap:10px;margin:5px 0}}
.bar{{height:18px;background:#2f5496;border-radius:3px;min-width:3px}}
.tab{{padding:7px 18px;cursor:pointer;background:#dce6f1;border:none;font-size:13px;border-radius:4px 4px 0 0;margin-right:3px}}
.tab.on{{background:#2f5496;color:#fff}}
.tp{{display:none}} .tp.on{{display:block}}
.ctx{{font-size:11px;color:#666;max-width:380px}}
</style>
</head>
<body>
<div class="header">
  <h1>&#128274; 개인정보/민감정보 스캔 결과</h1>
  <p>스캔 경로: {he(summary['scan_path'])} &nbsp;|&nbsp; 실행: {he(summary['scan_time'])}</p>
</div>
<div class="container">
<div class="cards">
  <div class="card"><div class="val">{summary['total_files']:,}</div><div class="lbl">스캔 파일</div></div>
  <div class="card red"><div class="val">{summary['files_with_pii']:,}</div><div class="lbl">개인정보 파일</div></div>
  <div class="card red"><div class="val">{summary['total_pii']:,}</div><div class="lbl">총 탐지 건수</div></div>
  <div class="card"><div class="val">{int(summary['elapsed'])//60}분 {int(summary['elapsed'])%60}초</div><div class="lbl">소요 시간</div></div>
</div>
<div style="margin-bottom:0">
  <button class="tab on" onclick="tab(0)">유형별 통계</button>
  <button class="tab" onclick="tab(1)">파일 목록</button>
  <button class="tab" onclick="tab(2)">상세 결과</button>
</div>

<div class="tp on" id="t0"><div class="section">
  <div class="sh">유형별 탐지 건수</div>
  <div style="padding:16px">
"""
        for t, c in type_counts.items():
            if c == 0: continue
            bw = max(4, c * 350 // max_count)
            html += f'    <div class="bar-row"><div style="width:100px;text-align:right;font-size:12px">{he(t)}</div>'
            html += f'<div class="bar" style="width:{bw}px"></div><div style="font-size:12px">{c}건</div></div>\n'
        html += "  </div></div></div>\n"

        html += '<div class="tp" id="t1"><div class="section"><div class="sh">파일 목록</div>\n'
        html += '<table><tr><th>파일 경로</th><th>확장자</th><th>탐지</th><th>추출 방법</th><th>소요(초)</th><th>오류</th></tr>\n'
        for r in results:
            cls = "red" if r.match_count > 0 else "gray"
            html += (f'<tr><td class="{cls}">{he(r.path)}</td>'
                     f'<td>{he(r.extension)}</td>'
                     f'<td class="{cls}">{r.match_count}</td>'
                     f'<td>{he(r.method)}</td>'
                     f'<td>{r.scan_sec:.2f}</td>'
                     f'<td style="color:#c00">{he(r.error)}</td></tr>\n')
        html += "</table></div></div>\n"

        html += '<div class="tp" id="t2"><div class="section"><div class="sh">상세 탐지 결과</div>\n'
        html += '<table><tr><th>파일</th><th>유형</th><th>탐지 값</th><th>마스킹</th><th>줄</th><th>맥락</th></tr>\n'
        for r in results:
            for m in r.matches:
                bc = badge_cls.get(m.pii_type, "other")
                html += (f'<tr><td>{he(r.path)}</td>'
                         f'<td><span class="badge b-{bc}">{he(m.pii_type)}</span></td>'
                         f'<td class="red">{he(m.matched)}</td>'
                         f'<td style="color:#888">{he(m.masked)}</td>'
                         f'<td style="text-align:center">{m.line_no}</td>'
                         f'<td class="ctx">{he(m.context)}</td></tr>\n')
        html += "</table></div></div>\n"

        html += """</div>
<script>
function tab(i){
  document.querySelectorAll('.tp').forEach((e,j)=>e.classList.toggle('on',i===j));
  document.querySelectorAll('.tab').forEach((e,j)=>e.classList.toggle('on',i===j));
}
</script>
</body></html>"""

        with open(path, "w", encoding="utf-8-sig") as f:
            f.write(html)


# ════════════════════════════════════════════════════════════════════
# 5. 메인 오케스트레이션
# ════════════════════════════════════════════════════════════════════

def parse_args():
    p = argparse.ArgumentParser(
        description="개인정보/민감정보 파일 스캐너",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""기본 동작:
  기본 스캔 경로 : 현재 사용자 폴더(%USERPROFILE%, %APPDATA% 등)
  시스템 폴더   : Windows/Program Files/ProgramData 등 항상 제외
  레지스트리    : 미사용 (환경변수로만 경로 조회)
  네트워크      : 스캔 중 미사용 (로컬 파일 읽기만 수행)

예시:
  pii_scanner.py                              # 내 사용자 폴더 전체
  pii_scanner.py --path "D:\\업무자료"         # 특정 폴더
  pii_scanner.py --all-drives                 # 전체 드라이브 (주의)
  pii_scanner.py --exclude "C:\\Projects"     # 특정 폴더 추가 제외
"""
    )
    p.add_argument("--path",        default="",  metavar="경로",
                   help="스캔할 특정 폴더 (미지정 시 사용자 폴더 자동 선택)")
    p.add_argument("--output",      default="reports", metavar="경로",
                   help="리포트 저장 폴더 (기본: ./reports)")
    p.add_argument("--dll",         default="",
                   help="Everything64.dll 경로 (미지정 시 자동 탐색)")
    p.add_argument("--skip-images", action="store_true",
                   help="이미지 파일 OCR 건너뜀 (속도 향상)")
    p.add_argument("--max-size",    type=int, default=100, metavar="MB",
                   help="최대 파일 크기 MB (기본 100)")
    p.add_argument("--threads",     type=int, default=0,
                   help="병렬 스레드 수 (기본 자동: CPU 코어 수)")
    p.add_argument("--all-drives",  action="store_true",
                   help="전체 드라이브 스캔 (시스템 폴더 제외는 유지)")
    p.add_argument("--exclude",     action="append", default=[], metavar="경로",
                   help="추가로 제외할 폴더 (여러 번 사용 가능)")
    p.add_argument("--include-system", action="store_true",
                   help="시스템 폴더 제외 비활성화 (권장하지 않음)")
    return p.parse_args()


def main():
    args = parse_args()

    print()
    print(f"{Fore.CYAN}{'='*60}")
    print(f"  PiiScanner - 개인정보/민감정보 스캐너 (Python)")
    print(f"{'='*60}{Style.RESET_ALL}")
    print(f"  {Fore.GREEN}[보안] 레지스트리 미접근 / 스캔 중 네트워크 미사용{Style.RESET_ALL}")
    print(f"  {Fore.GREEN}[보안] Windows 시스템 폴더 자동 제외{Style.RESET_ALL}")
    print()

    max_bytes = args.max_size * 1024 * 1024
    n_threads = args.threads or max(1, (os.cpu_count() or 4))
    exclude_system = not args.include_system

    # ── 스캔 경로 결정 ───────────────────────────────────────────
    if args.path:
        # 사용자가 명시한 경로
        root_paths = [args.path]
        scan_label = args.path
    elif args.all_drives:
        # 전체 드라이브 (명시적 opt-in)
        root_paths = []   # Everything에 경로 제한 없음
        scan_label = "전체 드라이브 (시스템 폴더 제외)"
        print(f"  {Fore.YELLOW}[주의] 전체 드라이브 스캔 - 시간이 오래 걸릴 수 있습니다.{Style.RESET_ALL}")
    else:
        # 기본: 현재 사용자 폴더만
        root_paths = get_user_default_paths()
        if not root_paths:
            root_paths = [os.environ.get("USERPROFILE", "C:\\Users")]
        scan_label = "사용자 폴더 (%USERPROFILE% 등)"
        print(f"  스캔 경로:")
        for p in root_paths:
            print(f"    {Fore.CYAN}{p}{Style.RESET_ALL}")
    print()

    # ── Step 1: Everything SDK 초기화 ───────────────────────────
    print(f"{Fore.YELLOW}[1/4] Everything SDK 초기화 중...{Style.RESET_ALL}")
    scanner = EverythingScanner()
    if not scanner.initialize(args.dll):
        print(f"\n{Fore.RED}Everything을 실행한 후 다시 시도하세요.")
        print(f"다운로드: https://www.voidtools.com/{Style.RESET_ALL}")
        sys.exit(1)
    print(f"  {Fore.GREEN}✓ Everything SDK 연결 성공{Style.RESET_ALL}")

    # ── Step 2: 파일 목록 조회 ───────────────────────────────────
    print(f"\n{Fore.YELLOW}[2/4] 파일 목록 조회 중 ({scan_label})...{Style.RESET_ALL}")
    if exclude_system:
        print(f"  (Windows/Program Files/ProgramData 등 시스템 폴더 자동 제외)")
    t0 = time.time()
    files = scanner.scan(
        root_paths=root_paths,
        exclude_system=exclude_system,
        extra_excludes=args.exclude,
    )
    t1 = time.time()
    doc_cnt = sum(1 for f in files if f.is_document)
    img_cnt = sum(1 for f in files if f.is_image)
    print(f"  {Fore.GREEN}✓ {len(files):,}개 발견 ({t1-t0:.1f}초){Style.RESET_ALL}")
    print(f"     문서: {doc_cnt:,}개  /  이미지: {img_cnt:,}개")

    # ── Step 3: PII 스캔 ─────────────────────────────────────────
    print(f"\n{Fore.YELLOW}[3/4] 개인정보 스캔 중 (스레드: {n_threads}개)...{Style.RESET_ALL}")
    detector = PiiDetector()
    results: List[FileScanResult] = []
    pii_total = 0
    lock = threading.Lock()

    def scan_file(entry: FileEntry) -> FileScanResult:
        if entry.size > max_bytes:
            return FileScanResult(entry.path, entry.extension, False,
                                  f"파일 크기 초과 ({entry.size//1024//1024}MB)",
                                  "skipped", 0, 0.0)
        t_s = time.time()
        text, method, error = extract_text(entry.path, entry.extension,
                                           use_ocr=(not args.skip_images))
        matches = detector.detect(text) if text else []
        t_e = time.time()
        return FileScanResult(
            path=entry.path, extension=entry.extension,
            success=bool(text), error=error, method=method,
            text_len=len(text), scan_sec=t_e-t_s, matches=matches
        )

    with tqdm(total=len(files), unit="파일", dynamic_ncols=True,
              bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]") as pbar:
        with ThreadPoolExecutor(max_workers=n_threads) as executor:
            futures = {executor.submit(scan_file, f): f for f in files}
            for future in as_completed(futures):
                res = future.result()
                with lock:
                    results.append(res)
                    pii_total += res.match_count
                pbar.update(1)
                pbar.set_postfix({"탐지": pii_total})

    elapsed = time.time() - t0
    files_with_pii = sum(1 for r in results if r.match_count > 0)
    print(f"\n  {Fore.GREEN}✓ 스캔 완료 | 탐지: {pii_total:,}건 | "
          f"소요: {int(elapsed)//60}분 {int(elapsed)%60}초{Style.RESET_ALL}")

    # ── Step 4: 리포트 생성 ──────────────────────────────────────
    print(f"\n{Fore.YELLOW}[4/4] 리포트 생성 중...{Style.RESET_ALL}")

    # 유형별 집계
    type_counts: Dict[str, int] = {}
    for r in results:
        for m in r.matches:
            type_counts[m.pii_type] = type_counts.get(m.pii_type, 0) + 1

    summary = {
        "scan_path":    scan_label,
        "scan_time":    datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "total_files":  len(files),
        "files_with_pii": files_with_pii,
        "total_pii":    pii_total,
        "elapsed":      elapsed,
        "type_counts":  dict(sorted(type_counts.items(),
                                    key=lambda x: x[1], reverse=True)),
    }

    # 결과 정렬: 탐지 많은 순
    results.sort(key=lambda r: r.match_count, reverse=True)

    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    reporter = Reporter()
    xlsx_path, html_path = reporter.save_all(
        results, summary, args.output, f"pii_report_{ts}"
    )

    print(f"  {Fore.GREEN}✓ Excel: {xlsx_path}")
    print(f"  ✓ HTML:  {html_path}{Style.RESET_ALL}")

    # ── 최종 출력 ────────────────────────────────────────────────
    print()
    print(f"{Fore.CYAN}{'='*52}")
    print("  스캔 결과 요약")
    print(f"{'='*52}{Style.RESET_ALL}")
    print(f"  총 파일 수:        {len(files):,}")
    print(f"  개인정보 파일 수:  {Fore.RED}{files_with_pii:,}{Style.RESET_ALL}")
    print(f"  총 탐지 건수:      {Fore.RED}{pii_total:,}{Style.RESET_ALL}")
    if type_counts:
        print()
        print("  [유형별]")
        for t, c in sorted(type_counts.items(), key=lambda x: x[1], reverse=True):
            print(f"    {t:<14}: {c:,}건")
    print(f"\n  소요: {int(elapsed)//60}분 {int(elapsed)%60}초")
    print(f"{Fore.CYAN}{'='*52}{Style.RESET_ALL}")
    print()


if __name__ == "__main__":
    main()
